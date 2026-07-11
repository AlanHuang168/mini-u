# -*- coding: utf-8 -*-
"""
生成零基础教学 PPT：Mini-U 量化分析员
运行： python build_ppt.py  → 输出 Mini-U教学课件.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 配色 ----------
NAVY   = RGBColor(0x1E, 0x3A, 0x8A)   # 深蓝 标题
CYAN   = RGBColor(0x08, 0x91, 0xB2)   # 青色 副标题
RED    = RGBColor(0xBE, 0x12, 0x3D)   # 强调红
DARK   = RGBColor(0x1F, 0x29, 0x37)   # 正文深灰
GRAY   = RGBColor(0x64, 0x74, 0x8B)   # 次要灰
TIPBG  = RGBColor(0xEC, 0xFE, 0xFF)   # 提示框底(青)
WARNBG = RGBColor(0xFE, 0xF2, 0xF2)   # 警告框底(红)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
         space_after=8, first=False, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def rounded_box(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(2)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_bar(slide, color=CYAN):
    """顶部装饰条"""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False


def content_slide(title, bullets, note="", tip=None, tip_warn=False):
    """通用内容页：标题 + 要点列表 + 可选提示框"""
    s = add_slide()
    set_bg(s, WHITE)
    add_bar(s)
    # 标题
    tf = textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(1))
    para(tf, title, 34, NAVY, bold=True, first=True)
    # 要点
    body = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.3))
    for i, (txt, size, color, bold) in enumerate(bullets):
        para(body, txt, size, color, bold=bold, first=(i == 0), space_after=12)
    # 提示框
    if tip:
        box_top = Inches(6.0)
        rounded_box(s, Inches(0.9), box_top, Inches(11.5), Inches(1.05),
                    WARNBG if tip_warn else TIPBG,
                    line=(RED if tip_warn else CYAN))
        ttf = textbox(s, Inches(1.2), box_top + Inches(0.1), Inches(11), Inches(0.85),
                      anchor=MSO_ANCHOR.MIDDLE)
        para(ttf, tip, 19, DARK, bold=False, first=True)
    if note:
        notes(s, note)
    return s


def table_slide(title, headers, rows, note="", tip=None):
    """带表格的页面"""
    s = add_slide()
    set_bg(s, WHITE)
    add_bar(s)
    tf = textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(1))
    para(tf, title, 34, NAVY, bold=True, first=True)

    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_w = Inches(11.6)
    tbl_h = Inches(0.6) * nrows
    left = Inches(0.85)
    top = Inches(1.9)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, tbl_w, tbl_h).table

    # 表头
    for c, h in enumerate(headers):
        cell = gtbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    # 数据行
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(16); r.font.color.rgb = DARK; r.font.name = FONT
            if c == 0:
                r.font.bold = True; r.font.color.rgb = CYAN

    if tip:
        box_top = Inches(6.0)
        rounded_box(s, Inches(0.9), box_top, Inches(11.5), Inches(1.0), TIPBG, line=CYAN)
        ttf = textbox(s, Inches(1.2), box_top + Inches(0.08), Inches(11), Inches(0.85),
                      anchor=MSO_ANCHOR.MIDDLE)
        para(ttf, tip, 18, DARK, first=True)
    if note:
        notes(s, note)
    return s


def code_slide(title, intro, code, note="", tip=None):
    """带代码块的页面"""
    s = add_slide(); set_bg(s, WHITE); add_bar(s)
    tf = textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(1))
    para(tf, title, 34, NAVY, bold=True, first=True)
    if intro:
        itf = textbox(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.8))
        para(itf, intro, 20, DARK, first=True)
    # 代码块底
    code_top = Inches(2.45)
    code_h = Inches(2.9)
    rounded_box(s, Inches(0.9), code_top, Inches(11.5), code_h, RGBColor(0x0F, 0x17, 0x2A))
    ctf = textbox(s, Inches(1.2), code_top + Inches(0.15), Inches(11), code_h - Inches(0.3))
    for i, line in enumerate(code.split("\n")):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = line if line else " "
        r.font.size = Pt(16); r.font.name = "Consolas"
        r.font.color.rgb = RGBColor(0x9C, 0xDC, 0xFE)
    if tip:
        box_top = Inches(5.7)
        rounded_box(s, Inches(0.9), box_top, Inches(11.5), Inches(1.1), TIPBG, line=CYAN)
        ttf = textbox(s, Inches(1.2), box_top + Inches(0.1), Inches(11), Inches(0.9),
                      anchor=MSO_ANCHOR.MIDDLE)
        para(ttf, tip, 18, DARK, first=True)
    if note:
        notes(s, note)
    return s


# ================= 幻灯片内容 =================

# --- 封面 ---
s = add_slide()
set_bg(s, NAVY)
from pptx.enum.shapes import MSO_SHAPE
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.7), SW, Inches(2.1))
band.fill.solid(); band.fill.fore_color.rgb = CYAN; band.line.fill.background(); band.shadow.inherit = False
tf = textbox(s, Inches(1), Inches(1.2), Inches(11.3), Inches(1.3))
para(tf, "🚀 Mini-U 量化分析员", 48, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf2 = textbox(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.5), anchor=MSO_ANCHOR.MIDDLE)
para(tf2, "用一个真实项目，讲透「AI 应用」是怎么做出来的", 30, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf3 = textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(1))
para(tf3, "零基础 · AI 应用开发入门课", 24, WHITE, align=PP_ALIGN.CENTER, first=True)
notes(s, "同学们好。今天我们不讲空洞的理论，而是拆解一个真实的、已经上线的 AI 项目——它能输入一个股票代码，自动生成一份 AI 分析报告。学完这节课，你会明白：所谓「AI 应用」，到底是由哪些零件拼起来的。")

# --- 本课收获 ---
content_slide(
    "这节课你会学到什么",
    [
        ("🧩  一个 AI 应用是由哪些零件组成的", 24, DARK, False),
        ("🔗  什么是「流水线（Pipeline）」思维", 24, DARK, False),
        ("🤝  规则和大模型（AI）如何分工合作", 24, DARK, False),
        ("💰  为什么要做「缓存」省钱", 24, DARK, False),
        ("🛡️  真实项目里踩过的坑（最宝贵的经验）", 24, DARK, False),
    ],
    tip="不需要任何基础 —— 只要你会用电脑、听说过「ChatGPT」，就能跟上。",
    note="这节课的目标不是让你马上会写代码，而是建立「AI 应用是怎么搭起来的」整体认知。我会用大量生活化的比喻，遇到专业词我都会翻译成人话。",
)

# --- 成果展示 ---
content_slide(
    "先看成果：它能干什么？",
    [
        ("你输入：一个股票代码，比如 600519（贵州茅台）", 24, RED, True),
        ("它输出：一份完整的分析报告 👇", 24, DARK, True),
        ("📊  股价、涨跌、综合评分", 22, DARK, False),
        ("🤖  一段 AI 写的专业分析", 22, DARK, False),
        ("👥  4 位「投资大师」的不同点评", 22, DARK, False),
        ("📄  一个漂亮的网页报告", 22, DARK, False),
    ],
    tip="一句话：输入代码 → 自动产出报告，全程几十秒。",
    note="先给大家看看最终效果。你只要输入一个股票代码，点一下按钮，系统就自动帮你采集数据、算分、写分析、请4位风格不同的大师点评，最后生成一份网页报告。这就是我们要拆解的目标。",
)

# --- 预备知识 表格 ---
table_slide(
    "预备知识：三个词，先听懂",
    ["词", "人话解释", "生活类比"],
    [
        ["大模型 / LLM", "会聊天、会写文章的 AI（如 DeepSeek）", "一个博学的「实习生」"],
        ["API", "程序之间打电话的「电话号码」", "点外卖的下单接口"],
        ["API Key", "打这个电话要用的「密码」", "你的外卖账号密码"],
    ],
    tip="⚠️ 记住：API Key 是要花钱的密码，泄露了别人就能花你的钱！",
    note="开始前先扫清三个术语。大模型就是像 ChatGPT 那样会写字的 AI，可以想象成一个很博学但需要你给指令的实习生。API 是程序之间互相调用的接口。API Key 是调用时的密码，而且是要花钱的——这点后面会专门讲一个真实的安全事故。",
)

# --- 核心思想：拆解 ---
code_slide(
    "核心思想：把大事拆成小步",
    "想象你要做一道菜 🍳，不会一步做完，而是：",
    "做菜：   买菜  →  切菜  →  下锅炒  →  装盘\n\nAI应用： 采集数据  →   AI分析   →  生成报告\n           (买菜)       (炒菜)        (装盘)",
    tip="每一步只干一件事 —— 出了问题，一眼就知道是哪一步坏了。",
    note="这是整节课最重要的思想。做复杂的事情，不要想着一步登天，而要拆成一个个简单的小步骤，像流水线一样串起来。做菜要买菜切菜炒菜装盘，我们的程序就是采集数据、AI分析、生成报告。每一步单独简单，串起来就完成了复杂任务。",
)

# --- 系统分层 ---
code_slide(
    "系统全貌：像一栋楼分层",
    "整个系统像一栋楼，分成四层：",
    "┌─────────────────────────────┐\n│  界面层：用户点点点的网页    │ ← 你看到的\n│  功能层：单股分析 / 多股对比 │ ← 不同功能\n│  流水线：采集 → 分析 → 报告  │ ← 干活核心\n│  底层：  数据源 / AI / 缓存  │ ← 基础设施\n└─────────────────────────────┘",
    tip="分层的好处：换了顶楼的装修，不影响地基。改界面不用动核心逻辑。",
    note="整个系统像一栋楼，分成几层。最上面是用户看到的网页界面，中间是各种功能，再下面是干活的流水线，最底下是数据源和AI这些基础设施。分层的最大好处是互不干扰——我想换个好看的界面，完全不用碰下面算分析的代码。",
)

# --- 第1步 采集打分 ---
code_slide(
    "第 1 步：采集数据 + 打分",
    "做两件事：① 抓行情（股价涨跌）  ② 用规则打分（满分100）",
    "打分规则（人人都看得懂）：\n\n   股价 > 1000 元   →  加 15 分\n   今天涨了         →  加 10 分\n   ROE > 20%        →  加 20 分",
    tip="规则写在配置文件里，改分数不用改程序 —— 运营同学也能维护。",
    note="流水线第一步，采集数据。程序自动去网上抓股价、涨跌，然后用一套规则打分。股价高就加分、上涨就加分。关键是——这些规则不是写死在代码里，而是放在一个人人能读懂的配置文件里，改分数不用懂编程。",
)

# --- 为什么要规则 ---
table_slide(
    "为什么要「规则」而不全交给 AI？",
    ["谁", "负责什么", "类比"],
    [
        ["规则", "算分数（精确的活）", "计算器"],
        ["AI", "写分析（表达的活）", "作家"],
    ],
    tip="AI 不擅长精确算数，还可能「一本正经地胡说」（叫幻觉）；规则算得准、能解释。",
    note="有同学会问，都有AI了为什么还要写规则？因为AI有两个毛病：一是不擅长精确算数，二是会幻觉——一本正经地编造。而规则引擎算得准、能解释清楚为什么打这个分。所以我们让它们分工。",
)

# --- 金句页 ---
s = add_slide()
set_bg(s, NAVY)
tf = textbox(s, Inches(1), Inches(1.0), Inches(11.3), Inches(1))
para(tf, "🌟 本课最重要的一句话", 28, CYAN, bold=True, align=PP_ALIGN.CENTER, first=True)
tf2 = textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
para(tf2, "让规则引擎做数学题", 44, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf2, "让大模型做作文题", 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
tf3 = textbox(s, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.3), anchor=MSO_ANCHOR.MIDDLE)
para(tf3, "把精确计算交给代码，把语言表达交给 AI —— 这是做可靠 AI 应用的核心原则。",
     20, CYAN, align=PP_ALIGN.CENTER, first=True)
notes(s, "这句话请大家记下来，它是做AI应用的黄金法则：让规则引擎做数学题，让大模型做作文题。精确的、不能出错的部分交给代码；需要灵活表达、写得好看的部分交给AI。这样既发挥了AI的长处，又避开了它的短处。")

# --- 第2步 AI分析 ---
code_slide(
    "第 2 步：AI 分析",
    "把数据和分数打包，交给大模型 DeepSeek：",
    "你是一位资深股票分析师。\n这是数据：{茅台，1190元，评分65…}\n请写一份 250-400 字的专业分析，\n最后给出评级：强烈推荐 / 推荐 / 观望 / 回避",
    tip="给 AI 的这段指令叫「Prompt（提示词）」—— 写好它是一门手艺。",
    note="第二步，把数据和分数交给大模型DeepSeek。我们给它一段指令，告诉它你是分析师、这是数据、请写一份分析。这段指令有个专业名字叫Prompt，也就是提示词。同样的数据，Prompt写得好不好，AI输出的质量天差地别。",
)

# --- 提示词技巧 ---
content_slide(
    "提示词的小技巧",
    [
        ("我们的提示词里明确要求：", 24, DARK, True),
        ("✅  结构：基本面 → 风险 → 逻辑 → 建议", 23, DARK, False),
        ("✅  字数：250-400 字（不能太长太短）", 23, DARK, False),
        ("✅  给选项：评级只能从「强烈推荐/推荐/观望/回避」里选", 23, DARK, False),
    ],
    tip="越具体，AI 越听话。别只说「帮我分析」，要说清楚「结构、长度、格式」。",
    note="写提示词有个诀窍：越具体越好。你不能只说帮我分析一下，那AI会天马行空。你要告诉它结构、字数、评级选项。把要求框得越清楚，AI的输出就越稳定可控。这个技巧大家平时用ChatGPT也能马上用上。",
)

# --- 第3步 多角色 ---
table_slide(
    "第 3 步：多位大师来点评",
    ["大师", "风格", "关注点"],
    [
        ["巴菲特", "价值投资", "护城河、安全边际"],
        ["赵老哥", "游资短线", "板块热度、情绪"],
        ["木头姐", "成长创新", "颠覆性、增长"],
        ["陈小群", "宏观策略", "政策、行业周期"],
    ],
    tip="每个大师就是一张「人设卡」。加新大师 = 加一张卡，不用改程序。",
    note="第三步很有意思。同一只股票，我们让AI分别扮演巴菲特、赵老哥这些风格迥异的投资大师来点评，得到多个视角。每个大师就是一张人设卡，是个几行的小文件。你想加新大师，只要复制一张卡改改内容，完全不用碰程序代码。",
)

# --- 缓存 ---
code_slide(
    "省钱小能手：缓存",
    "每次调用 AI 都要花钱、还要等几秒。同一只股票反复分析，太浪费！",
    "有人问茅台\n  → 先看「小本本」上有没有记过\n     记过且没过期  →  直接念出来（不花钱！）\n     没记过        →  真的问 AI，然后记到本本",
    tip="缓存有「保质期」（本项目 6 小时）—— 因为股价会变，太旧的不能用。",
    note="调用AI是要花钱的，而且要等待。如果同一只股票几分钟内被问了十次，每次都真去问AI，就是白白烧钱。所以我们加了缓存——像个小本本，问过的答案先记下来，下次直接念。但缓存有保质期，因为股价会变，太旧的答案就不能用了，我们设的是6小时。",
)

# --- 事故1 密码泄露 ---
content_slide(
    "🛡️ 真实事故 1：密码泄露",
    [
        ("这个项目上线后，发现一个安全问题：", 24, DARK, True),
        ("用户 A 输入了自己的 API Key（花钱的密码），", 23, RED, False),
        ("结果用户 B 打开网页，也能用 A 的密码！= 花 A 的钱 😱", 23, RED, True),
        ("", 12, DARK, False),
        ("原因：程序错误地把密码存在了「所有人共享」的地方。", 23, DARK, False),
        ("修复：改成「每个人一个独立小房间」，互相看不见。", 23, DARK, False),
    ],
    tip="⚠️ 敏感信息（密码、密钥）绝不能存在「大家共享」的地方。",
    tip_warn=True,
    note="现在讲最宝贵的部分——真实踩过的坑。这个项目上线后发现：用户A输入了自己的付费密码，结果用户B打开网页居然也能用，等于花A的钱。这是很严重的安全事故。原因是程序把密码存到了所有人共享的区域。修复方法是给每个用户一个独立的小房间，互相隔离。",
)

# --- 事故1教训 ---
content_slide(
    "事故 1 的教训",
    [
        ("敏感信息绝不能存在「大家共享」的地方。", 26, RED, True),
        ("", 12, DARK, False),
        ("❌  把银行卡密码写在公司公告栏 → 人人可见", 23, DARK, False),
        ("✅  把密码放进自己的保险箱 → 只有你能开", 23, DARK, False),
    ],
    tip="做多用户的网站，一定要问自己：「这个数据，别的用户会不会看到？」",
    note="这个事故的教训是：密码这类敏感信息，绝不能存在大家共享的地方。就像你不会把银行卡密码贴在公司公告栏上。做任何多用户的网站，都要养成习惯问自己：我存的这个东西，别的用户会不会看到？",
)

# --- 事故2 崩溃 ---
content_slide(
    "🛡️ 真实事故 2：神秘崩溃",
    [
        ("现象：同一份代码——", 24, DARK, True),
        ("在服务器上（Linux）✅ 正常", 23, DARK, False),
        ("在电脑命令行 ✅ 正常", 23, DARK, False),
        ("但用网页跑 ❌ 直接崩溃", 23, RED, True),
        ("", 12, DARK, False),
        ("最后定位：某个底层组件在特定环境下不兼容 → 加「演示模式」绕过。", 22, DARK, False),
    ],
    tip="排查思路：对比找不同（系统？运行方式？），把可疑部分单独拎出来测。",
    note="第二个坑更玄乎。同样的代码，服务器上正常，命令行正常，就是网页一跑就崩溃。核心方法是对比找不同、单独拎出来测。我们发现区别在运行方式，最后定位到一个底层组件在特定环境下不兼容。解决办法是加了个演示模式开关绕过它。这个排查思路你以后遇到任何bug都用得上。",
)

# --- 事故2教训 ---
content_slide(
    "事故 2 的教训",
    [
        ("「在我电脑上好好的呀」", 28, RED, True),
        ("—— 程序员最常说、也最坑的一句话。", 22, GRAY, False),
        ("", 12, DARK, False),
        ("同样代码换个环境就可能出问题，因为环境有差异：", 23, DARK, False),
        ("• 操作系统不同（Windows / Mac / Linux）", 22, DARK, False),
        ("• 软件版本不同 · 运行方式不同", 22, DARK, False),
    ],
    tip="出 bug 别慌：对比环境差异 + 逐步缩小范围。",
    note="这个坑的教训，也是程序员那句经典甩锅——在我电脑上好好的呀。同样代码换个环境就出问题太常见了，因为操作系统、软件版本、运行方式都可能不一样。遇到bug不要慌，冷静地对比环境差异，一步步缩小范围。",
)

# --- 事故3&4 ---
content_slide(
    "🛡️ 真实事故 3 & 4：小疏忽",
    [
        ("事故 3：忘了声明依赖", 24, RED, True),
        ("程序用到一个工具包，却忘了写进「购物清单」。", 22, DARK, False),
        ("👉 教训：程序需要什么，必须明明白白列出来。", 22, CYAN, False),
        ("", 10, DARK, False),
        ("事故 4：把巨大的临时文件传上去了", 24, RED, True),
        ("一不小心把几百 MB 临时文件提交到了代码仓库。", 22, DARK, False),
        ("👉 教训：临时文件、密钥文件不该进仓库，要写进「黑名单」。", 22, CYAN, False),
    ],
    note="还有两个小疏忽。一个是忘了声明依赖——程序要用某个工具包却没写进清单，换台电脑就找不到。教训是程序需要什么必须白纸黑字列清楚。另一个是不小心把几百MB临时文件传到了代码仓库。教训是临时文件和密钥文件不该进仓库，要提前写进黑名单。这些都是新手最常犯的错误。",
)

# --- 复盘全流程 ---
code_slide(
    "复盘：一次完整的分析之旅",
    "",
    "输入「600519」\n  ① 采集：抓行情 + 规则打分 → 65分\n  ② AI分析：把数据交给 DeepSeek → 写出报告\n  ③ 多角色：4位大师分别点评\n  ④ 生成漂亮的网页报告 + 保存\n  ✅ 你看到完整报告",
    tip="数据的旅程：数字 → 文字 → 网页。几乎所有 AI 应用都是这个套路。",
    note="我们把整个流程串一遍。你输入600519，第一步采集数据打分得65分，第二步交给AI写分析，第三步请4位大师点评，第四步拼成网页报告并保存。注意数据的形态变化：从数字，变成文字，最后变成网页。你以后看到的绝大多数AI应用都是这个套路。",
)

# --- 总结 ---
content_slide(
    "总结：Mini-U 教会我们 5 件事",
    [
        ("🧩  拆解 —— 把复杂任务拆成流水线的小步骤", 23, DARK, False),
        ("🤝  分工 —— 规则做数学题，AI 做作文题", 23, DARK, False),
        ("⚙️  配置驱动 —— 规则/角色/提示词外置，改配置不改代码", 23, DARK, False),
        ("💰  工程优化 —— 缓存省钱、出错兜底", 23, DARK, False),
        ("🛡️  安全意识 —— 密钥隔离、环境差异、别踩坑", 23, DARK, False),
    ],
    tip="AI 应用 ≠ 把问题丢给 AI，而是用工程手段把 AI 编排进一条可控的流水线。",
    note="最后总结，这个小项目教会我们5件事：学会拆解、懂得分工、配置驱动、工程优化、安全意识。请记住最后这句话：做AI应用不是简单地把问题丢给AI，而是用工程的手段，把AI这个聪明但不完全可靠的伙伴，编排进一条你能掌控的流水线里。",
)

# --- 练习 ---
table_slide(
    "课后动手练习",
    ["难度", "任务"],
    [
        ["⭐", "加一位新的投资大师（改一个配置文件）"],
        ["⭐", "在规则里加一条新的加分项"],
        ["⭐⭐", "把报告加一张评分图表"],
        ["⭐⭐⭐", "把 DeepSeek 换成另一个大模型"],
    ],
    tip="动手改一次，胜过看十遍。从加一位大师开始，最简单也最有成就感。",
    note="光听不练是学不会的。给大家留了几个练习，从易到难。强烈建议从加一位新大师开始——你只要复制一个配置文件，改改名字和风格，就能看到你的大师出现在报告里，特别有成就感。改一次代码，胜过看十遍PPT。",
)

# --- 结尾 ---
s = add_slide()
set_bg(s, NAVY)
tf = textbox(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.5), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "🎉 谢谢大家！", 48, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf2 = textbox(s, Inches(1), Inches(3.4), Inches(11.3), Inches(2), anchor=MSO_ANCHOR.MIDDLE)
para(tf2, "项目已开源 · 在线体验：HuggingFace Space  dayahs/mini-u", 22, CYAN, align=PP_ALIGN.CENTER, first=True)
para(tf2, "完整技术文档：见项目 TECH_OVERVIEW.md", 22, CYAN, align=PP_ALIGN.CENTER)
para(tf2, "有问题？欢迎提问 🙋", 24, WHITE, bold=True, align=PP_ALIGN.CENTER)
notes(s, "今天的课就到这里。整个项目已经开源，大家可以去HuggingFace上亲自体验，代码和更详细的技术文档都在项目里。希望这节课让你对AI应用到底怎么做有了具体的感觉。有任何问题欢迎现在提问或课后交流。谢谢大家！")

# ---------- 保存 ----------
out = "Mini-U教学课件.pptx"
prs.save(out)
print(f"✅ 已生成：{out}  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页")
